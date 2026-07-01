"""PPTX text extractor using python-pptx."""

import os

from . import TextExtractor, ExtractedText


class PptxExtractor(TextExtractor):
	"""Extractor for PPTX / PPT PowerPoint files."""

	def extract(self, file_path: str) -> ExtractedText:
		"""Extract text from a PowerPoint presentation."""
		try:
			from pptx import Presentation
		except ImportError:
			raise ImportError(
				"PPTX extraction requires 'python-pptx'. "
				"Install with: pip install python-pptx"
			)

		prs = Presentation(file_path)
		slide_count = len(prs.slides)
		text_parts = []

		for idx, slide in enumerate(prs.slides, start=1):
			slide_lines = [f"## Slide {idx}"]
			for shape in slide.shapes:
				if hasattr(shape, "text") and shape.text.strip():
					slide_lines.append(shape.text.strip())
			if len(slide_lines) > 1:
				text_parts.append("\n".join(slide_lines))

		return ExtractedText(
			text="\n\n".join(text_parts),
			title=os.path.basename(file_path),
			metadata={"file_type": "pptx", "slides": slide_count},
		)
